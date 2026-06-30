# -*- coding: utf-8 -*-
"""Offline presentation builder — generates professional .pptx files without WPS.
Requires: pip install python-pptx Pillow
"""
from typing import Optional, Dict, List
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "line": MSO_SHAPE.ISOSCELES_TRIANGLE,  # fallback
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
} if HAS_PPTX else {}

ALIGN_PPTX = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY} if HAS_PPTX else {}


def _hex_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def build_pptx(structure: Dict, output_path: str) -> Dict:
    """Build .pptx from JSON structure with design palette support."""
    if not HAS_PPTX:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    try:
        from intelligence.design_rules import PPTX_PALETTES
    except ImportError:
        PPTX_PALETTES = {}

    prs = Presentation()
    prs.slide_width = Inches(13.333) if structure.get("layout") == "16x9" else Inches(10)
    prs.slide_height = Inches(7.5) if structure.get("layout") == "16x9" else Inches(5.625)

    palette_key = structure.get("palette", "")
    palette = PPTX_PALETTES.get(palette_key, {})
    primary_color = palette.get("primary", "1E2761")
    secondary_color = palette.get("secondary", "CADCFC")
    accent_color = palette.get("accent", "FFFFFF")

    stats = {"slides": 0}

    for slide_data in structure.get("slides", []):
        layout_idx = 0  # default blank
        layout_name = slide_data.get("layout", "blank")
        for i, layout in enumerate(prs.slide_layouts):
            if layout_name.lower() in layout.name.lower():
                layout_idx = i
                break
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Background
        bg_color = slide_data.get("background", "")
        if bg_color:
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = _hex_rgb(bg_color)

        # Title
        title_cfg = slide_data.get("title", {})
        if title_cfg:
            title_shape = slide.shapes.title
            if title_shape:
                tf = title_shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title_cfg.get("text", "")
                p.font.size = Pt(title_cfg.get("size_pt", 36))
                p.font.bold = title_cfg.get("bold", True)
                p.font.color.rgb = _hex_rgb(title_cfg.get("color", accent_color))

        # Bullets / body text
        bullets = slide_data.get("bullets", [])
        if bullets:
            left = Inches(slide_data.get("left", 1.0))
            top = Inches(slide_data.get("top", 1.8))
            width = Inches(slide_data.get("width", 8.0))
            height = Inches(slide_data.get("height", 3.5))
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(slide_data.get("body_size_pt", 16))
                p.font.color.rgb = _hex_rgb(slide_data.get("body_color", primary_color))
                p.level = slide_data.get("indent_level", 0)

        # Image
        img_cfg = slide_data.get("image", {})
        if img_cfg and img_cfg.get("path"):
            path = img_cfg["path"]
            if Path(path).exists():
                slide.shapes.add_picture(
                    path,
                    Inches(img_cfg.get("left", 5.0)),
                    Inches(img_cfg.get("top", 1.5)),
                    Inches(img_cfg.get("width", 4.0)),
                    Inches(img_cfg.get("height", 3.0)),
                )

        # Shape
        shape_cfg = slide_data.get("shape", {})
        if shape_cfg:
            st = SHAPE_MAP.get(shape_cfg.get("type", "rectangle"), MSO_SHAPE.RECTANGLE)
            s = slide.shapes.add_shape(
                st,
                Inches(shape_cfg.get("left", 1.0)),
                Inches(shape_cfg.get("top", 2.0)),
                Inches(shape_cfg.get("width", 3.0)),
                Inches(shape_cfg.get("height", 2.0)),
            )
            if shape_cfg.get("fill_color"):
                s.fill.solid()
                s.fill.fore_color.rgb = _hex_rgb(shape_cfg["fill_color"])

        stats["slides"] += 1

    prs.save(output_path)
    return {"output": output_path, "stats": stats}


def extract_pptx_text(filepath: str) -> Dict:
    """Extract all text from a .pptx file."""
    if not HAS_PPTX:
        return {"error": "python-pptx not installed"}
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_text.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    slide_text.append(" | ".join(row_text))
        slides.append({"index": i, "text": slide_text})
    return {"slides": slides, "total_slides": len(slides)}
